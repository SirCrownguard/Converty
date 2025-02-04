import os, sys, zipfile, json, subprocess, threading, csv, datetime
try:
    import pdf2image
    import pptx
    import tqdm
    import colorama
    import tkinter as tk
except ImportError as e:
    from tkinter import messagebox
    root = tk.Tk()
    root.withdraw()
    messagebox.showerror("Missing Library", f"Required library missing: {str(e)}\nPlease use full EXE version or pip install {str(e).split()[-1]}")
    sys.exit(1)

from tkinter import filedialog, ttk
from pdf2image import convert_from_path
from pptx import Presentation
from tqdm import tqdm
from colorama import Fore, init
from tempfile import NamedTemporaryFile

init(autoreset=True)
try:
    import comtypes.client
except ImportError:
    pass

PREFERENCES_FILE = "preferences.json"
HISTORY_FILE = "history.csv"

COLORS = {
    "background": "#DAD7CD",
    "primary": "#A3B18A",
    "secondary": "#588157",
    "accent": "#3A5A40",
    "text": "#344E41",
    "success": "#6A994E",
    "warnings": "#B85C5C",
    "errors": "#9E2B25"
}

DARK_COLORS = {
    "background": "#344E41",
    "primary": "#3A5A40",
    "secondary": "#588157",
    "accent": "#A3B18A",
    "text": "#DAD7CD",
    "success": "#A7C957",
    "warnings": "#B85C5C",
    "errors": "#9E2B25"
}

FONTS = {
    "default": ("Open Sans", 11),
    "title": ("Open Sans", 13, "bold"),
    "button": ("Open Sans", 11, "bold")
}

LANGUAGES = {
    "en": {
        "select_theme": "Select Theme:",
        "default_theme": "Light",
        "dark_theme": "Dark",
        "select_language": "Select Language:",
        "select_conversion": "Select Conversion Type:",
        "pdf_to_pptx": "PDF to PPTX",
        "pptx_to_pdf": "PPTX to PDF",
        "select_mode": "Processing Mode:",
        "select_zip": "Save results as ZIP file",
        "ok": "OK",
        "select_pdf": "Select a PDF file",
        "select_multiple_pdfs": "Select multiple PDF files",
        "select_pptx": "Select a PPTX file",
        "select_multiple_pptxs": "Select multiple PPTX files",
        "select_folder": "Select a folder",
        "no_file_selected": "No files selected! Process cancelled.",
        "no_output_folder": "No output folder selected! Process cancelled.",
        "completed": "Process completed! Files saved:",
        "zip_completed": "Process completed! All files are in",
        "save_preferences": "Save Preferences",
        "select_pdf_engine": "Select PDF conversion engine:",
        "powerpoint_com": "PowerPoint COM",
        "libreoffice": "LibreOffice",
        "icon_error": "Error loading icon:",
        "progress_title": "Progress",
        "processing": "Processing: {}",
        "finished": "Process completed!",
        "close": "Close",
        "start_conversion": "Start Conversion",
        "operation_history": "Operation History:",
        "clear_history": "Clear History",
        "date": "Date",
        "conversion_type": "Conversion Type",
        "processed_files": "Processed Files",
        "output_location": "Output Location",
        "compressed": "Compressed",
        "mode": "Mode",
        "files_processed": "Files Processed",
        "settings": "Settings"
    },
    "tr": {
        "select_theme": "Tema Seçimi:",
        "default_theme": "Varsayılan",
        "dark_theme": "Koyu",
        "select_language": "Dil Seçimi:",
        "select_conversion": "Dönüşüm Türü Seçimi:",
        "pdf_to_pptx": "PDF'ten PPTX'e",
        "pptx_to_pdf": "PPTX'ten PDF'e",
        "select_mode": "İşlem Modu:",
        "select_zip": "Sonuçları ZIP olarak kaydet",
        "ok": "Tamam",
        "select_pdf": "PDF dosyanı seç",
        "select_multiple_pdfs": "Birden fazla PDF dosyası seç",
        "select_pptx": "PPTX dosyanı seç",
        "select_multiple_pptxs": "Birden fazla PPTX dosyası seç",
        "select_folder": "Bir klasör seç",
        "no_file_selected": "Seçili dosya yok! İşlem iptal edildi.",
        "no_output_folder": "Çıktı klasörü seçilmedi! İşlem iptal edildi.",
        "completed": "İşlem tamamlandı! Kaydedilen dosyalar:",
        "zip_completed": "İşlem tamamlandı! Tüm dosyalar:",
        "save_preferences": "Tercihleri Kaydet",
        "select_pdf_engine": "PDF dönüşüm motorunu seç:",
        "powerpoint_com": "PowerPoint COM (Sadece Windows)",
        "libreoffice": "LibreOffice",
        "icon_error": "İkon yüklenirken hata oluştu:",
        "progress_title": "İlerleme",
        "processing": "İşleniyor: {}",
        "finished": "İşlem tamamlandı!",
        "close": "Kapat",
        "start_conversion": "Dönüştür",
        "operation_history": "İşlem Geçmişi:",
        "clear_history": "Geçmişi Temizle",
        "date": "Tarih",
        "conversion_type": "Dönüşüm Türü",
        "processed_files": "İşlenen Dosyalar",
        "output_location": "Kayıt Konumu",
        "compressed": "Sıkıştırıldı",
        "mode": "Mod",
        "files_processed": "İşlenen Dosyalar",
        "settings": "Ayarlar"
    }
}

pdf_engine_options = {
    "powerpoint_com": {
         "en": LANGUAGES["en"]["powerpoint_com"],
         "tr": LANGUAGES["tr"]["powerpoint_com"]
    },
    "libreoffice": {
         "en": LANGUAGES["en"]["libreoffice"],
         "tr": LANGUAGES["tr"]["libreoffice"]
    }
}

def load_preferences():
    if os.path.exists(PREFERENCES_FILE):
        with open(PREFERENCES_FILE, "r") as file:
            return json.load(file)
    return {}

def save_preferences(preferences):
    with open(PREFERENCES_FILE, "w") as file:
        json.dump(preferences, file)

def ensure_history_file():
    if not os.path.exists(HISTORY_FILE):
        with open(HISTORY_FILE, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Date", "ConversionKey", "Mode", "Compressed", "Files", "OutputLocation"])

def clear_history_file():
    with open(HISTORY_FILE, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Date", "ConversionKey", "Mode", "Compressed", "Files", "OutputLocation"])

def add_history_entry(date_str, conversion_key, mode, compressed, files, output_location):
    ensure_history_file()
    with open(HISTORY_FILE, "a", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow([
            date_str,
            conversion_key,
            mode,
            "✓" if compressed else "✗",
            files,
            output_location
        ])

def load_history():
    ensure_history_file()
    rows = []
    with open(HISTORY_FILE, "r", newline="", encoding="utf-8") as f:
        reader = csv.reader(f)
        next(reader, None)
        for row in reader:
            rows.append(row)
    return rows

def get_pdf_files(mode, lang):
    root = tk.Tk()
    root.withdraw()
    if mode == 1:
        pdf_path = filedialog.askopenfilename(title=LANGUAGES[lang]["select_pdf"], filetypes=[("PDF Files", "*.pdf")])
        return [pdf_path] if pdf_path else []
    elif mode == 2:
        pdf_paths = filedialog.askopenfilenames(title=LANGUAGES[lang]["select_multiple_pdfs"], filetypes=[("PDF Files", "*.pdf")])
        return list(pdf_paths) if pdf_paths else []
    elif mode == 3:
        folder_path = filedialog.askdirectory(title=LANGUAGES[lang]["select_folder"])
        return [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.lower().endswith(".pdf")] if folder_path else []
    return []

def get_pptx_files(mode, lang):
    root = tk.Tk()
    root.withdraw()
    if mode == 1:
        pptx_path = filedialog.askopenfilename(title=LANGUAGES[lang]["select_pptx"], filetypes=[("PPTX Files", "*.pptx")])
        return [pptx_path] if pptx_path else []
    elif mode == 2:
        pptx_paths = filedialog.askopenfilenames(title=LANGUAGES[lang]["select_multiple_pptxs"], filetypes=[("PPTX Files", "*.pptx")])
        return list(pptx_paths) if pptx_paths else []
    elif mode == 3:
        folder_path = filedialog.askdirectory(title=LANGUAGES[lang]["select_folder"])
        return [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.lower().endswith(".pptx")] if folder_path else []
    return []

def pdf_to_pptx(lang, mode, zip_option, progress_update=None):
    pdf_files = get_pdf_files(mode, lang)
    if not pdf_files:
        print(Fore.RED + LANGUAGES[lang]["no_file_selected"])
        return []
    output_folder = filedialog.askdirectory(title=LANGUAGES[lang]["select_folder"])
    if not output_folder:
        print(Fore.RED + LANGUAGES[lang]["no_output_folder"])
        return []
    output_files = []
    total = len(pdf_files)
    for i, pdf_path in enumerate(pdf_files):
        if progress_update:
            progress_update(i+1, total, LANGUAGES[lang]["processing"].format(os.path.basename(pdf_path)))
        output_pptx_path = os.path.join(output_folder, os.path.basename(pdf_path).replace(".pdf", ".pptx"))
        images = convert_from_path(pdf_path)
        presentation = Presentation()
        for image in images:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            iw, ih = image.size
            sw, sh = presentation.slide_width, presentation.slide_height
            ratio = iw / ih
            if ratio > (sw / sh):
                nw = sw
                nh = sw / ratio
            else:
                nh = sh
                nw = sh * ratio
            left = (sw - nw) / 2
            top = (sh - nh) / 2
            with NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                image.save(tmp.name)
                slide.shapes.add_picture(tmp.name, left, top, width=nw, height=nh)
                temp_name = tmp.name
            os.remove(temp_name)
        presentation.save(output_pptx_path)
        output_files.append(output_pptx_path)
    if zip_option:
        zip_filename = os.path.join(output_folder, "converted_pptx_files.zip")
        with zipfile.ZipFile(zip_filename, "w") as zipf:
            for file in output_files:
                zipf.write(file, os.path.basename(file))
                os.remove(file)
        print(Fore.GREEN + f"\n{LANGUAGES[lang]['zip_completed']} {zip_filename}")
    else:
        print(Fore.GREEN + f"\n{LANGUAGES[lang]['completed']}")
        for out in output_files:
            print(Fore.CYAN + out)
    return output_files

def pptx_to_pdf(lang, mode, zip_option, pdf_engine, progress_update=None):
    pptx_files = get_pptx_files(mode, lang)
    if not pptx_files:
        print(Fore.RED + LANGUAGES[lang]["no_file_selected"])
        return []
    output_folder = filedialog.askdirectory(title=LANGUAGES[lang]["select_folder"])
    if not output_folder:
        print(Fore.RED + LANGUAGES[lang]["no_output_folder"])
        return []
    output_files = []
    total = len(pptx_files)
    if pdf_engine == "powerpoint_com":
        try:
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
        except Exception as e:
            print(Fore.RED + f"PowerPoint başlatılamadı: {e}")
            return []
        for i, pptx_path in enumerate(pptx_files):
            if progress_update:
                progress_update(i+1, total, LANGUAGES[lang]["processing"].format(os.path.basename(pptx_path)))
            pptx_path = os.path.abspath(pptx_path)
            if not os.path.exists(pptx_path):
                print(Fore.RED + f"Dosya bulunamadı: {pptx_path}")
                continue
            output_pdf_path = os.path.join(output_folder, os.path.basename(pptx_path).replace(".pptx", ".pdf"))
            try:
                presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
                presentation.ExportAsFixedFormat(output_pdf_path, 2, Intent=2)
                presentation.Close()
                output_files.append(output_pdf_path)
            except Exception as e:
                print(Fore.RED + f"{pptx_path} dönüştürülürken hata oluştu: {e}")
        powerpoint.Quit()
    elif pdf_engine == "libreoffice":
        for i, pptx_path in enumerate(pptx_files):
            if progress_update:
                progress_update(i+1, total, LANGUAGES[lang]["processing"].format(os.path.basename(pptx_path)))
            try:
                subprocess.run(["soffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", output_folder],
                               check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                output_pdf_path = os.path.join(output_folder, os.path.basename(pptx_path).replace(".pptx", ".pdf"))
                output_files.append(output_pdf_path)
            except subprocess.CalledProcessError as e:
                print(Fore.RED + f"{pptx_path} dönüştürülürken hata oluştu: {e}")
    else:
        print(Fore.RED + "Tanımlı PDF dönüşüm motoru bulunamadı!")
        return []
    if zip_option:
        zip_filename = os.path.join(output_folder, "converted_pdf_files.zip")
        with zipfile.ZipFile(zip_filename, "w") as zipf:
            for file in output_files:
                zipf.write(file, os.path.basename(file))
                os.remove(file)
        print(Fore.GREEN + f"\n{LANGUAGES[lang]['zip_completed']} {zip_filename}")
    else:
        print(Fore.GREEN + f"\n{LANGUAGES[lang]['completed']}")
        for out in output_files:
            print(Fore.CYAN + out)
    return output_files

def main_app():
    prefs = load_preferences()
    root = tk.Tk()
    root.title("Converty")
    root.geometry("1300x1300")
    style = ttk.Style(root)
    style.theme_use('clam')
    style.configure(".", background=COLORS["background"], foreground=COLORS["text"], font=FONTS["default"])
    # TButton için orijinal ayarları kaldırdık
    style.configure("TRadiobutton", background=COLORS["primary"], foreground=COLORS["text"], indicatorsize=16, indicatormargin=4, padding=5)
    style.configure("TCombobox", fieldbackground=COLORS["primary"], selectbackground=COLORS["accent"])
    style.configure("Treeview", background=COLORS["primary"], fieldbackground=COLORS["primary"], foreground=COLORS["text"], rowheight=25)
    style.configure("Treeview.Heading", background=COLORS["accent"], foreground=COLORS["primary"], font=FONTS["button"])
    style.configure("Horizontal.TProgressbar", thickness=20, troughcolor=COLORS["secondary"], troughrelief='flat', background=COLORS["success"])
    root.configure(bg=COLORS["background"])

    lang_var = tk.StringVar(value=prefs.get("language", "tr"))
    conv_type_var = tk.StringVar(value=prefs.get("conversion_type", "pdf_to_pptx"))
    mode_var = tk.IntVar(value=prefs.get("mode", 1))
    zip_var = tk.BooleanVar(value=prefs.get("zip_option", False))
    save_pref_var = tk.BooleanVar(value=False)
    pdf_engine_var_display = tk.StringVar()
    default_engine = prefs.get("pdf_engine", "libreoffice")
    pdf_engine_var_display.set(pdf_engine_options[default_engine][lang_var.get()])
    theme_var = tk.StringVar(value=prefs.get("theme", "Default"))

    main_frame = ttk.Frame(root)
    main_frame.pack(fill="both", expand=True, padx=20, pady=20)

    # Oluşturulacak widget'ların başlıklarını boş bırakıp update_labels ile dolduracağız.
    pref_frame = ttk.LabelFrame(main_frame, text="", padding=(15, 10))
    pref_frame.pack(fill="x", pady=10, padx=15)

    top_frame = ttk.Frame(pref_frame)
    top_frame.pack(fill="x", pady=5)
    theme_frame = ttk.Frame(top_frame)
    theme_frame.pack(side="left", padx=5)
    # Tema label'i update_labels'de oluşturulacak
    theme_label = ttk.Label(theme_frame, text="")
    theme_label.pack(anchor="w")
    theme_combo = ttk.Combobox(theme_frame, textvariable=theme_var, state="readonly", width=10)
    theme_combo.pack(anchor="w", padx=5)

    lang_frame = ttk.Frame(top_frame)
    lang_frame.pack(side="left", padx=20)
    # Dil label'i update_labels'de oluşturulacak
    lang_label = ttk.Label(lang_frame, text="")
    lang_label.pack(anchor="w")
    lang_combo = ttk.Combobox(lang_frame, textvariable=lang_var, values=list(LANGUAGES.keys()), state="readonly", width=10)
    lang_combo.pack(anchor="w", padx=5)

    conv_subframe = ttk.Frame(pref_frame)
    conv_subframe.pack(fill="x", padx=5, pady=5)
    # Dönüşüm label'i update_labels'de oluşturulacak
    conv_label = ttk.Label(conv_subframe, text="")
    conv_label.pack(anchor="w")
    conv_radio1 = ttk.Radiobutton(conv_subframe, text="", variable=conv_type_var, value="pdf_to_pptx")
    conv_radio1.pack(anchor="w")
    conv_radio2 = ttk.Radiobutton(conv_subframe, text="", variable=conv_type_var, value="pptx_to_pdf")
    conv_radio2.pack(anchor="w")

    # İşlem modu label'i update_labels'de oluşturulacak
    mode_label = ttk.Label(pref_frame, text="")
    mode_label.pack(anchor="w", padx=5, pady=5)
    mode_frame = ttk.Frame(pref_frame)
    mode_frame.pack(anchor="w", padx=5, pady=5)

    zip_check = ttk.Checkbutton(pref_frame, text="", variable=zip_var)
    zip_check.pack(anchor="w", padx=5, pady=5)

    pdf_eng_frame = ttk.Frame(pref_frame)
    pdf_eng_frame.pack(anchor="w", padx=5, pady=5)
    pdf_eng_label = ttk.Label(pdf_eng_frame, text="")
    pdf_eng_label.pack(side="left", padx=5)
    pdf_engine_combo = ttk.Combobox(pdf_eng_frame, textvariable=pdf_engine_var_display, state="readonly", width=20)
    pdf_engine_combo.pack(side="left", padx=5)

    save_pref_check = ttk.Checkbutton(pref_frame, text="", variable=save_pref_var)
    save_pref_check.pack(anchor="w", padx=5, pady=5)

    convert_btn = ttk.Button(main_frame, text="")
    convert_btn.pack(pady=10)

    progress_frame = ttk.LabelFrame(main_frame, text="")
    progress_frame.pack(fill="x", pady=10)
    progress_label = ttk.Label(progress_frame, text="")
    progress_label.pack(pady=5)
    progress_bar = ttk.Progressbar(progress_frame, length=400, mode="determinate")
    progress_bar.pack(pady=5)

    history_frame = ttk.LabelFrame(main_frame, text="")
    history_frame.pack(fill="both", expand=True, pady=10)
    tree = ttk.Treeview(history_frame, columns=("date", "conversion", "mode", "zip", "files", "output"), show="headings")
    vsb = ttk.Scrollbar(history_frame, orient="vertical", command=tree.yview)
    tree.configure(yscrollcommand=vsb.set)
    tree.grid(row=0, column=0, sticky="nsew")
    vsb.grid(row=0, column=1, sticky="ns")
    history_frame.grid_columnconfigure(0, weight=1)
    history_frame.grid_rowconfigure(0, weight=1)
    clear_btn = ttk.Button(history_frame, text="")
    clear_btn.grid(row=1, column=0, pady=5, sticky="e")

    def update_style():
        if theme_var.get().lower() in ["koyu", "dark"]:
            current_colors = DARK_COLORS
        else:
            current_colors = COLORS
        style.configure(".", background=current_colors["background"], foreground=current_colors["text"])
        style.configure("TButton", background=current_colors["accent"], foreground=current_colors["primary"])
        style.configure("TRadiobutton", background=current_colors["primary"], foreground=current_colors["text"])
        style.configure("TCombobox", fieldbackground=current_colors["primary"], selectbackground=current_colors["accent"])
        style.configure("Treeview", background=current_colors["primary"], fieldbackground=current_colors["primary"], foreground=current_colors["text"])
        style.configure("Treeview.Heading", background=current_colors["accent"], foreground=current_colors["primary"])
        style.configure("Horizontal.TProgressbar", troughcolor=current_colors["secondary"], background=current_colors["success"])
        root.configure(bg=current_colors["background"])

    def refresh_history():
        for item in tree.get_children():
            tree.delete(item)
        current_lang = lang_var.get()
        for row in load_history():
            date_str, conversion_key, mode_num, compressed, files, output = row
            try:
                translated_conversion = LANGUAGES[current_lang][conversion_key]
            except KeyError:
                translated_conversion = conversion_key
            try:
                mode_int = int(mode_num)
            except ValueError:
                mode_int = 1
            if conversion_key == "pdf_to_pptx":
                mode_options = ["select_pdf", "select_multiple_pdfs", "select_folder"]
            else:
                mode_options = ["select_pptx", "select_multiple_pptxs", "select_folder"]
            if 1 <= mode_int <= 3:
                mode_desc = LANGUAGES[current_lang][mode_options[mode_int - 1]]
            else:
                mode_desc = "Unknown"
            tree.insert("", "end", values=(date_str, translated_conversion, mode_desc, compressed, files, output))

    def clear_history():
        clear_history_file()
        refresh_history()

    def update_labels(*args):
        lang = lang_var.get()
        # Update frame ve label başlıkları
        pref_frame.config(text=LANGUAGES[lang]["settings"])
        theme_label.config(text=LANGUAGES[lang]["select_theme"])
        lang_label.config(text=LANGUAGES[lang]["select_language"])
        conv_label.config(text=LANGUAGES[lang]["select_conversion"])
        conv_radio1.config(text=LANGUAGES[lang]["pdf_to_pptx"])
        conv_radio2.config(text=LANGUAGES[lang]["pptx_to_pdf"])
        mode_label.config(text=LANGUAGES[lang]["select_mode"])
        zip_check.config(text=LANGUAGES[lang]["select_zip"])
        pdf_eng_label.config(text=LANGUAGES[lang]["select_pdf_engine"])
        save_pref_check.config(text=LANGUAGES[lang]["save_preferences"])
        convert_btn.config(text=LANGUAGES[lang]["start_conversion"])
        progress_frame.config(text=LANGUAGES[lang]["progress_title"])
        clear_btn.config(text=LANGUAGES[lang]["clear_history"])
        tree.heading("date", text=LANGUAGES[lang]["date"])
        tree.heading("conversion", text=LANGUAGES[lang]["conversion_type"])
        tree.heading("mode", text=LANGUAGES[lang]["mode"])
        tree.heading("zip", text=LANGUAGES[lang]["compressed"])
        tree.heading("files", text=LANGUAGES[lang]["processed_files"])
        tree.heading("output", text=LANGUAGES[lang]["output_location"])

        # Güncel dil için tema combobox değerleri
        if lang == "tr":
            theme_combo['values'] = [LANGUAGES["tr"]["default_theme"], LANGUAGES["tr"]["dark_theme"]]
        else:
            theme_combo['values'] = [LANGUAGES["en"]["default_theme"], LANGUAGES["en"]["dark_theme"]]

        # PDF Engine seçenekleri
        pdf_engine_combo['values'] = [
            pdf_engine_options["powerpoint_com"][lang],
            pdf_engine_options["libreoffice"][lang]
        ]
        # Mode radio button'ları yeniden oluşturuluyor
        for widget in mode_frame.winfo_children():
            widget.destroy()
        if conv_type_var.get() == "pdf_to_pptx":
            texts = [
                LANGUAGES[lang]["select_pdf"],
                LANGUAGES[lang]["select_multiple_pdfs"],
                LANGUAGES[lang]["select_folder"]
            ]
        else:
            texts = [
                LANGUAGES[lang]["select_pptx"],
                LANGUAGES[lang]["select_multiple_pptxs"],
                LANGUAGES[lang]["select_folder"]
            ]
        for i, text in enumerate(texts):
            ttk.Radiobutton(mode_frame, text=text, variable=mode_var, value=i+1).pack(anchor="w", fill="x")
        refresh_history()
        update_style()

    def start_conversion():
        convert_btn.config(state="disabled")
        progress_bar["value"] = 0
        progress_label.config(text="")
        lang = lang_var.get()
        conversion = conv_type_var.get()
        mode = mode_var.get()
        zip_option = zip_var.get()
        pdf_engine = None
        if conversion == "pptx_to_pdf":
            selected_display = pdf_engine_var_display.get()
            for key, translations in pdf_engine_options.items():
                if translations[lang] == selected_display:
                    pdf_engine = key
                    break
        if save_pref_var.get():
            new_prefs = {
                "language": lang,
                "conversion_type": conversion,
                "mode": mode,
                "zip_option": zip_option,
                "pdf_engine": pdf_engine,
                "theme": theme_var.get()
            }
            save_preferences(new_prefs)
        def update_progress(current, total, message):
            progress_bar["maximum"] = total
            progress_bar["value"] = current
            progress_label.config(text=message)
            if current == total:
                progress_label.config(text=LANGUAGES[lang]["finished"])
        def run_conv():
            try:
                if conversion == "pdf_to_pptx":
                    result = pdf_to_pptx(lang, mode, zip_option, lambda c, t, m: root.after(0, lambda: update_progress(c, t, m)))
                else:
                    result = pptx_to_pdf(lang, mode, zip_option, pdf_engine, lambda c, t, m: root.after(0, lambda: update_progress(c, t, m)))
                now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                output_location = os.path.dirname(result[0]) if result else ""
                conversion_key = "pdf_to_pptx" if conversion == "pdf_to_pptx" else "pptx_to_pdf"
                add_history_entry(now, conversion_key, mode, zip_option, len(result) if result else 0, output_location)
                refresh_history()
            finally:
                root.after(0, lambda: convert_btn.config(state="normal"))
        threading.Thread(target=run_conv, daemon=True).start()

    lang_var.trace_add("write", update_labels)
    conv_type_var.trace_add("write", update_labels)
    theme_var.trace_add("write", update_labels)
    convert_btn.configure(command=start_conversion)
    clear_btn.configure(command=clear_history)
    update_labels()
    refresh_history()
    root.mainloop()

if __name__ == "__main__":
    main_app()
